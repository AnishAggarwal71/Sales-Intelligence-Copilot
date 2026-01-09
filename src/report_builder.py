"""
Report Builder Module
Generates PowerPoint presentations with charts and insights
"""

import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path
import io
import sys
import logging

# Import settings
sys.path.append(str(Path(__file__).parent.parent))
from config.settings import (
    PPT_TITLE,
    PPT_SUBTITLE,
    CHART_WIDTH_INCHES,
    CHART_HEIGHT_INCHES,
    REPORTS_DIR,
    ASSETS_DIR
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportBuilder:
    """Build PowerPoint reports with charts and insights"""
    
    def __init__(self):
        """Initialize report builder"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str = PPT_TITLE, subtitle: str = PPT_SUBTITLE):
        """
        Add title slide
        
        Args:
            title: Presentation title
            subtitle: Presentation subtitle
        """
        logger.info("Adding title slide...")
        
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"{subtitle}\n{datetime.now().strftime('%B %d, %Y')}"
    
    def add_text_slide(self, title: str, content: str):
        """
        Add text slide with title and bullet points
        
        Args:
            title: Slide title
            content: Text content
        """
        logger.info(f"Adding text slide: {title}")
        
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # Add content (handle bullet points)
        for line in content.split('\n'):
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.level = 0
                p.font.size = Pt(14)
    
    def add_chart_slide(self, title: str, chart_image_bytes: bytes):
        """
        Add slide with embedded chart
        
        Args:
            title: Slide title
            chart_image_bytes: Chart image as bytes
        """
        logger.info(f"Adding chart slide: {title}")
        
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = title
        
        # Format title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)
        
        # Add chart image
        left = Inches(1)
        top = Inches(1.5)
        
        # Save bytes to temp file and add to slide
        img_stream = io.BytesIO(chart_image_bytes)
        slide.shapes.add_picture(img_stream, left, top, width=Inches(8))
    
    def add_metrics_slide(self, title: str, metrics: dict):
        """
        Add slide with key metrics
        
        Args:
            title: Slide title
            metrics: Dictionary of metrics to display
        """
        logger.info(f"Adding metrics slide: {title}")
        
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)
        
        # Add metrics in a grid
        start_top = Inches(1.5)
        start_left = Inches(1)
        box_width = Inches(3.5)
        box_height = Inches(1.2)
        spacing = Inches(0.3)
        
        metric_items = list(metrics.items())
        for idx, (key, value) in enumerate(metric_items):
            row = idx // 2
            col = idx % 2
            
            left = start_left + col * (box_width + spacing)
            top = start_top + row * (box_height + spacing)
            
            # Add metric box
            box = slide.shapes.add_textbox(left, top, box_width, box_height)
            tf = box.text_frame
            tf.word_wrap = True
            
            # Metric label
            p1 = tf.paragraphs[0]
            p1.text = key
            p1.font.size = Pt(14)
            p1.font.color.rgb = RGBColor(100, 100, 100)
            
            # Metric value
            p2 = tf.add_paragraph()
            p2.text = str(value)
            p2.font.size = Pt(24)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(31, 119, 180)
            
            # Add border (using shape)
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, box_width, box_height
            )
            shape.fill.background()
            shape.line.color.rgb = RGBColor(200, 200, 200)
            shape.line.width = Pt(1)
    
    def add_table_slide(self, title: str, df: pd.DataFrame, max_rows: int = 10):
        """
        Add slide with data table
        
        Args:
            title: Slide title
            df: DataFrame to display
            max_rows: Maximum rows to show
        """
        logger.info(f"Adding table slide: {title}")
        
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)
        
        # Prepare data
        df_display = df.head(max_rows).copy()
        
        # Add table
        rows, cols = len(df_display) + 1, len(df_display.columns)  # +1 for header
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5) * (rows + 1)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths (must be integers in EMUs)
        col_width = int(width / cols)
        for col_idx in range(cols):
            table.columns[col_idx].width = col_width
        
        # Add header row
        for col_idx, col_name in enumerate(df_display.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Add data rows
        for row_idx, row in enumerate(df_display.itertuples(index=False), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    def save(self, filepath: str = None) -> str:
        """
        Save presentation to file
        
        Args:
            filepath: Path to save file
            
        Returns:
            Path where file was saved
        """
        if filepath is None:
            REPORTS_DIR.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filepath = REPORTS_DIR / f"sales_intelligence_report_{timestamp}.pptx"
        else:
            filepath = Path(filepath)
        
        filepath.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(filepath))
        
        logger.info(f"✓ Report saved to {filepath}")
        return str(filepath)


def create_chart_image(fig: go.Figure, width: int = 1200, height: int = 600) -> bytes:
    """
    Convert Plotly figure to image bytes
    
    Args:
        fig: Plotly figure
        width: Image width in pixels
        height: Image height in pixels
        
    Returns:
        Image as bytes
    """
    img_bytes = fig.to_image(format="png", width=width, height=height)
    return img_bytes


def build_full_report(metrics: dict, forecast_df: pd.DataFrame, 
                     at_risk_df: pd.DataFrame, insights: dict,
                     forecaster=None) -> str:
    """
    Build complete PowerPoint report
    
    Args:
        metrics: Business metrics dictionary
        forecast_df: Forecast DataFrame
        at_risk_df: At-risk customers DataFrame
        insights: Dictionary with executive summary and recommendations
        forecaster: Forecaster object (for getting charts)
        
    Returns:
        Path to saved report
    """
    logger.info("=" * 60)
    logger.info("BUILDING POWERPOINT REPORT")
    logger.info("=" * 60)
    
    builder = ReportBuilder()
    
    # Slide 1: Title
    builder.add_title_slide()
    
    # Slide 2: Executive Summary
    builder.add_text_slide(
        "Executive Summary",
        insights.get('executive_summary', 'Summary not available')
    )
    
    # Slide 3: Key Metrics
    key_metrics = {
        'Current MRR': f"${metrics.get('current_mrr', 0):,.0f}",
        'ARR': f"${metrics.get('arr', 0):,.0f}",
        'Active Customers': f"{metrics.get('active_customers', 0):,}",
        'Churn Rate': f"{metrics.get('churn_rate', 0):.1f}%",
        'ARPU': f"${metrics.get('arpu', 0):,.2f}",
        'CLV': f"${metrics.get('clv', 0):,.2f}"
    }
    builder.add_metrics_slide("Key Performance Indicators", key_metrics)
    
    # Slide 4: MRR Forecast Chart
    if forecaster:
        try:
            forecast_summary = forecaster.get_forecast_summary()
            
            # Create forecast chart
            fig = go.Figure()
            
            # Actual values
            actual_data = forecast_summary[forecast_summary['actual'].notna()]
            fig.add_trace(go.Scatter(
                x=actual_data['date'],
                y=actual_data['actual'],
                mode='lines',
                name='Actual MRR',
                line=dict(color='#1f77b4', width=3)
            ))
            
            # Predicted values
            fig.add_trace(go.Scatter(
                x=forecast_summary['date'],
                y=forecast_summary['predicted'],
                mode='lines',
                name='Predicted MRR',
                line=dict(color='#ff7f0e', width=3, dash='dash')
            ))
            
            # Confidence interval
            fig.add_trace(go.Scatter(
                x=forecast_summary['date'],
                y=forecast_summary['upper_bound'],
                mode='lines',
                line=dict(width=0),
                showlegend=False
            ))
            
            fig.add_trace(go.Scatter(
                x=forecast_summary['date'],
                y=forecast_summary['lower_bound'],
                mode='lines',
                fill='tonexty',
                fillcolor='rgba(255, 127, 14, 0.2)',
                line=dict(width=0),
                name='Confidence Interval'
            ))
            
            fig.update_layout(
                title="MRR Forecast - Next 90 Days",
                xaxis_title="Date",
                yaxis_title="MRR ($)",
                font=dict(size=14),
                showlegend=True,
                height=600
            )
            
            chart_bytes = create_chart_image(fig)
            builder.add_chart_slide("Revenue Forecast", chart_bytes)
        except Exception as e:
            logger.warning(f"Could not add forecast chart: {e}")
    
    # Slide 5: At-Risk Customers
    if at_risk_df is not None and len(at_risk_df) > 0:
        # Format the dataframe for display
        display_df = at_risk_df.head(10).copy()
        display_df['churn_probability'] = (display_df['churn_probability'] * 100).round(1).astype(str) + '%'
        display_df['current_revenue'] = display_df['current_revenue'].apply(lambda x: f"${x:.2f}")
        display_df['risk_score'] = display_df['risk_score'].round(2)
        
        # Select columns to display
        display_cols = ['customer_id', 'churn_probability', 'current_revenue', 'risk_score']
        if 'tenure_months' in display_df.columns:
            display_df['tenure_months'] = display_df['tenure_months'].round(1)
            display_cols.append('tenure_months')
        
        builder.add_table_slide(
            "Top At-Risk Customers",
            display_df[display_cols]
        )
    
    # Slide 6: Recommendations
    builder.add_text_slide(
        "Strategic Recommendations",
        insights.get('recommendations', 'Recommendations not available')
    )
    
    # Save report
    filepath = builder.save()
    
    logger.info("✓ Report generation complete!")
    return filepath


# Test the module
if __name__ == "__main__":
    import sys
    from pathlib import Path
    
    sys.path.append(str(Path(__file__).parent.parent))
    
    # Create a simple test report
    test_metrics = {
        'current_mrr': 245892.50,
        'arr': 2950710.00,
        'active_customers': 8234,
        'churn_rate': 5.21,
        'arpu': 29.87,
        'clv': 1847.32
    }
    
    test_insights = {
        'executive_summary': """
• Current MRR stands at $245,892 with +2.3% month-over-month growth
• 8,234 active customers with average revenue per user of $29.87
• Churn rate at 5.21% with 15 high-value customers identified as at-risk
        """.strip(),
        'recommendations': """
1. Deploy targeted retention campaigns for at-risk customers
2. Assign dedicated CSMs to high-value accounts
3. Analyze usage patterns to identify early churn signals
        """.strip()
    }
    
    # Create simple forecast data
    import pandas as pd
    dates = pd.date_range('2024-01-01', periods=12, freq='M')
    test_forecast = pd.DataFrame({
        'date': dates,
        'actual': [240000 + i*5000 for i in range(12)],
        'predicted': [242000 + i*5000 for i in range(12)]
    })
    
    # Create at-risk customers data
    test_at_risk = pd.DataFrame({
        'customer_id': [f'CUST_{i:04d}' for i in range(1, 11)],
        'churn_probability': [0.89, 0.87, 0.82, 0.78, 0.76, 0.74, 0.72, 0.71, 0.70, 0.69],
        'current_revenue': [499.0, 499.0, 299.0, 299.0, 99.0, 99.0, 99.0, 99.0, 29.0, 29.0],
        'risk_score': [445.11, 434.13, 245.18, 233.22, 75.24, 73.26, 71.28, 70.29, 20.30, 20.01],
        'tenure_months': [24.5, 18.2, 36.1, 12.8, 42.3, 8.5, 28.9, 15.2, 6.1, 48.7]
    })
    
    print("Building test report...")
    
    filepath = build_full_report(
        metrics=test_metrics,
        forecast_df=test_forecast,
        at_risk_df=test_at_risk,
        insights=test_insights,
        forecaster=None
    )
    
    print(f"\n✓ Test report created: {filepath}")
    print("Open the file to view the presentation!")